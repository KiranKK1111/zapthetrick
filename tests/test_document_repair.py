"""In-place document repair & beautify (app.documents.repair)."""
from __future__ import annotations

import asyncio
import io

from app.documents import repair as R


# --------------------------------------------------------------------------- #
# Markdown (pure, no libs)
# --------------------------------------------------------------------------- #
def test_markdown_heading_and_lists():
    out, changes = R.beautify_markdown("#Title\ntext\n* a\n+ b\n\n\n\nmore")
    assert out.startswith("# Title")
    assert "\n- a\n" in out and "\n- b\n" in out          # * / + → -
    assert "\n\n\n" not in out                            # blanks collapsed


def test_markdown_unbalanced_fence_warned():
    _, changes = R.beautify_markdown("```python\ncode without close")
    assert any("unbalanced" in c for c in changes)


def test_markdown_never_raises_on_garbage():
    out, _ = R.beautify_markdown("\x00\x01 weird")
    assert isinstance(out, str)


# --------------------------------------------------------------------------- #
# Excel
# --------------------------------------------------------------------------- #
def _xlsx_bytes():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Item", "Qty", "Price"])
    ws.append(["A", 2, 10])
    ws.append(["B", 3, 20])
    b = io.BytesIO()
    wb.save(b)
    return b.getvalue()


def test_xlsx_beautify_and_automation():
    from openpyxl import load_workbook
    res = R.repair_xlsx(_xlsx_bytes(), automate=True)
    assert res.ok and res.changed
    wb = load_workbook(io.BytesIO(res.data))
    ws = wb.active
    assert ws["A1"].font.bold is True                     # header bold
    assert ws.freeze_panes == "A2"                        # header frozen
    assert ws.auto_filter.ref is not None                 # auto-filter
    assert ws.cell(row=4, column=1).value == "TOTAL"      # SUM row
    assert str(ws.cell(row=4, column=2).value).startswith("=SUM(")
    assert ws.cell(row=5, column=1).value == "AVERAGE"    # AVERAGE row
    assert str(ws.cell(row=5, column=2).value).startswith("=AVERAGE(")
    assert len(ws._charts) == 1                            # bar chart added


def test_xlsx_beautify_without_automation():
    from openpyxl import load_workbook
    res = R.repair_xlsx(_xlsx_bytes(), automate=False)
    assert res.ok
    ws = load_workbook(io.BytesIO(res.data)).active
    assert ws.cell(row=4, column=1).value is None          # no totals row


def test_xlsx_bad_bytes_fail_open():
    res = R.repair_xlsx(b"not an xlsx")
    assert res.ok is False and "failed" in res.reason


# --------------------------------------------------------------------------- #
# Word / PowerPoint / PDF (round-trip validity)
# --------------------------------------------------------------------------- #
def test_docx_roundtrip():
    import docx
    d = docx.Document()
    for t in ["Hi ", "", "", "Bye"]:
        d.add_paragraph(t)
    b = io.BytesIO()
    d.save(b)
    res = R.repair_docx(b.getvalue())
    assert res.ok
    docx.Document(io.BytesIO(res.data))                    # reloads without error


def test_pptx_roundtrip():
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5]).shapes.title.text = "T"
    b = io.BytesIO()
    prs.save(b)
    res = R.repair_pptx(b.getvalue())
    assert res.ok
    Presentation(io.BytesIO(res.data))                     # reloads


def test_pdf_roundtrip():
    import fitz
    doc = fitz.open()
    doc.new_page().insert_text((72, 72), "hi")
    pb = doc.tobytes()
    doc.close()
    res = R.repair_pdf(pb)
    assert res.ok
    d2 = fitz.open(stream=res.data, filetype="pdf")
    assert d2.page_count == 1
    d2.close()


def test_pdf_detects_scanned_page():
    import fitz
    doc = fitz.open()
    pg = doc.new_page()
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 100, 100))
    pix.clear_with(200)
    pg.insert_image(fitz.Rect(0, 0, 100, 100), pixmap=pix)   # image, no text
    pb = doc.tobytes()
    doc.close()
    res = R.repair_pdf(pb)
    assert res.ok
    # OCR either ran (Tesseract present) or was reported as needed.
    assert any("OCR" in c or "scanned" in c.lower() for c in res.changes)


# --------------------------------------------------------------------------- #
# Dispatcher + code
# --------------------------------------------------------------------------- #
def test_dispatch_by_extension():
    async def _run():
        res = await R.repair_document(_xlsx_bytes(), "sheet.xlsx", automate=True)
        assert res.ok and res.kind == "xlsx" and res.data
        md = await R.repair_document(b"#Hi\ntext", "notes.md")
        assert md.kind == "md" and md.text.startswith("# Hi")
        code = await R.repair_document(b"x=1", "s.py")
        assert code.kind == "code"
        unknown = await R.repair_document(b"plain", "file.unknownext")
        assert unknown.ok and unknown.kind == "text"
    asyncio.run(_run())


def test_repair_code_never_raises():
    res = asyncio.run(R.repair_code("def f( ):\n  return 1", "python"))
    assert res.ok and res.kind == "code"
    assert isinstance(res.lint, list)


def test_lint_fix_fail_open_without_toolchain():
    # No ruff/eslint installed → fix_code returns the source unchanged (never raises).
    from app.polyglot.linters import fix_code, lint_code
    src = "import os\nx=1"
    assert asyncio.run(fix_code("python", src)) in (src, )  # unchanged when no tool
    assert asyncio.run(lint_code("cobol", "x")) == []       # unknown language
    assert asyncio.run(fix_code("python", "")) == ""
