"""P2-8 — PowerPoint (.pptx) document export (doc-skill parity).

Generates real .pptx bytes via render_document and validates them by reopening
with python-pptx: slide count, the title slide, a content slide's bullets, a
table slide, and a code slide. Plus the doc-request detection for pptx.
"""
from __future__ import annotations

import io

import pytest

from app.documents.detect import explicit_doc_formats, explicit_doc_request
from app.documents.generators import media_type, normalize_format, render_document

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402

DECK_MD = """\
# Caching Strategies

A short deck on caching.

## Why cache

- Reduce latency
- Offload the database
- Smooth traffic spikes

## Trade-offs

| Approach | Pro | Con |
| --- | --- | --- |
| Write-through | consistent | slower writes |
| Write-back | fast writes | risk on crash |

## Example

```python
cache.set("k", value, ttl=60)
```
"""


def _open(data: bytes) -> Presentation:
    return Presentation(io.BytesIO(data))


def test_pptx_format_registry():
    assert normalize_format("powerpoint") == "pptx"
    assert media_type("pptx") == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation")


def test_render_pptx_is_valid_and_multi_slide():
    data, mt, ext = render_document(DECK_MD, "pptx", "Caching 101")
    assert ext == "pptx"
    assert data[:2] == b"PK"      # .pptx is a zip
    prs = _open(data)
    slides = list(prs.slides)
    # title slide + Why cache + Trade-offs(table) + Example(code) ≥ 4
    assert len(slides) >= 4


def test_pptx_title_and_bullets_present():
    data, _, _ = render_document(DECK_MD, "pptx", "Caching 101")
    prs = _open(data)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
    blob = "\n".join(texts)
    assert "Caching 101" in blob            # title slide
    assert "Why cache" in blob              # section heading → slide title
    assert "Reduce latency" in blob         # a bullet
    assert "Write-through" in blob          # table cell
    assert "cache.set" in blob              # code slide


def test_pptx_has_a_table_shape():
    data, _, _ = render_document(DECK_MD, "pptx", "Deck")
    prs = _open(data)
    assert any(shape.has_table for slide in prs.slides
               for shape in slide.shapes)


def test_pptx_without_title_still_renders():
    data, _, _ = render_document("# Only\n\n- a\n- b\n", "pptx", "")
    prs = _open(data)
    assert len(list(prs.slides)) >= 1


def test_pptx_empty_content_no_crash():
    data, _, _ = render_document("", "pptx", "Empty")
    prs = _open(data)
    assert len(list(prs.slides)) >= 1       # at least the title slide


def test_long_bullet_list_spills_to_multiple_slides():
    md = "## Big\n\n" + "\n".join(f"- item {i}" for i in range(30))
    data, _, _ = render_document(md, "pptx", "")
    prs = _open(data)
    titles = [s.shapes.title.text for s in prs.slides
              if s.shapes.title is not None]
    assert any("(cont.)" in t for t in titles)


# ── detection ────────────────────────────────────────────────────────────────
def test_detect_powerpoint_requests():
    assert explicit_doc_request("give me a powerpoint")[1] == "pptx"
    assert explicit_doc_request("make a slide deck about caching")[1] == "pptx"
    assert explicit_doc_request("summarize this as a presentation")[1] == "pptx"


def test_detect_does_not_false_fire_on_it_terms():
    # "presentation layer" is an architecture term, not a doc request.
    assert explicit_doc_request(
        "I need a presentation layer for my app")[0] is False
    assert explicit_doc_request(
        "how do I open a pptx in python")[0] is False
