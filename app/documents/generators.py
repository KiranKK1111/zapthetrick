"""Generate downloadable documents from Markdown/text content.

Claude-style document export: the assistant's answer (or one of its
artifacts) is Markdown, and the user picks a format to download —
Markdown, plain text, CSV, Excel, Word, or PDF. This module converts the
Markdown into the requested format and returns the raw bytes.

  render_document(content, fmt, title="") -> (bytes, media_type, ext)

`content` is Markdown (or already-tabular text). For the spreadsheet
formats (csv/xlsx) we parse the GFM tables out of the content; for the
document formats (docx/pdf) we render the Markdown blocks (headings,
paragraphs, lists, code, quotes, tables) with light styling.
"""
from __future__ import annotations

import csv as _csv
import io
import os
import re

# pip names: python-docx -> `docx`, fpdf2 -> `fpdf`, openpyxl -> `openpyxl`.

# Bundled Unicode fonts (DejaVu) so the PDF can render smart dashes, arrows,
# non-breaking hyphens, box-drawing in ASCII banners, etc. — the built-in
# Helvetica/Courier are Latin-1 only and turn those into '?'.
_FONT_DIR = os.path.join(os.path.dirname(__file__), "fonts")

# Canonical formats and their HTTP media types.
_MEDIA = {
    "txt": "text/plain; charset=utf-8",
    "md": "text/markdown; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "csv": "text/csv; charset=utf-8",
    "json": "application/json; charset=utf-8",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
    "zip": "application/zip",
    "7z": "application/x-7z-compressed",
}

# Friendly synonyms the UI / model might send.
_ALIASES = {
    "markdown": "md",
    "mkd": "md",
    "htm": "html",
    "webpage": "html",
    "text": "txt",
    "plain": "txt",
    "word": "docx",
    "msword": "docx",
    "doc": "docx",
    "excel": "xlsx",
    "spreadsheet": "xlsx",
    "xls": "xlsx",
    "sheet": "xlsx",
    "powerpoint": "pptx",
    "ppt": "pptx",
    "slides": "pptx",
    "slide deck": "pptx",
    "deck": "pptx",
    "presentation": "pptx",
    "slideshow": "pptx",
    "geojson": "json",
    "jsonl": "json",
    "7zip": "7z",
    "7-zip": "7z",
    "sevenz": "7z",
    "p7zip": "7z",
}

SUPPORTED_FORMATS = tuple(_MEDIA.keys())


class UnsupportedFormat(ValueError):
    """Requested an export format we don't generate."""


def normalize_format(fmt: str) -> str:
    f = (fmt or "").strip().lower().lstrip(".")
    f = _ALIASES.get(f, f)
    if f in _MEDIA:
        return f
    # Source-code file: a language name / extension, or the generic "code"
    # token (language resolved from the answer at render time).
    from .detect import _CODE_EXTS, _CODE_LANG

    if f == "code":
        return "code"
    if f in _CODE_LANG:
        return _CODE_LANG[f]
    if f in _CODE_EXTS:
        return f
    raise UnsupportedFormat(
        f"Unsupported format '{fmt}'. Use one of: "
        + ", ".join(SUPPORTED_FORMATS)
    )


def media_type(fmt: str) -> str:
    # Code/source files are served as UTF-8 text; the extension drives the rest.
    return _MEDIA.get(normalize_format(fmt), "text/plain; charset=utf-8")


def _apply_branding_md(content: str, es) -> str:
    """Inject branding that Markdown CAN carry (header / confidentiality / author
    text + footer) so the PROSE renderers (PDF/DOCX/PPTX/TXT/MD) show it too.
    Accent color + logo image remain HTML-only (they need renderer-native APIs).
    No-op when there are no settings."""
    if es is None:
        return content
    top: list[str] = []
    if getattr(es, "confidentiality", ""):
        top.append(f"> {es.confidentiality}")
    if getattr(es, "header", ""):
        top.append(f"**{es.header}**")
    if getattr(es, "author", ""):
        top.append(f"*{es.author}*")
    head = ("\n\n".join(top) + "\n\n") if top else ""
    foot = (f"\n\n---\n\n*{es.footer}*\n") if getattr(es, "footer", "") else ""
    return head + content + foot


def _enrich_prose_markdown(content: str, title: str, language: str = "") -> str:
    """Phase-4 structure (TOC / numbering / glossary / appendix / auto-diagram)
    for the PROSE formats, applied by round-tripping through the DocumentModel —
    so PDF/DOCX/PPTX pick up the SAME enrichment as HTML. Gated by
    `auto_structure`; OFF → the content is returned unchanged (byte-identical
    legacy output). ``language`` (Phase 7) localizes the furniture labels.
    Fail-open to the raw content."""
    try:
        from .structure import auto_structure_enabled, enrich
        if not auto_structure_enabled():
            return content
        from .model import markdown_to_model, model_to_markdown
        return model_to_markdown(enrich(markdown_to_model(content, title),
                                        lang=language))
    except Exception:  # noqa: BLE001
        return content


def _model_driven_enabled() -> bool:
    """Phase 1: render the binary formats (PDF/DOCX/PPTX) straight from the
    DocumentModel IR. Default ON; `cfg.documents.model_driven_render = False`
    restores the legacy Markdown-tuple renderers."""
    try:
        from app.core.config_loader import cfg
        return bool(getattr(cfg.documents, "model_driven_render", True))
    except Exception:  # noqa: BLE001
        return True


def _prose_model(content: str, title: str, export_settings, language: str = ""):
    """Build the ONE DocumentModel the binary renderers consume (Phase 1): brand
    the Markdown, parse it into the IR, then apply the Phase-4 structure
    enrichment ON THE MODEL (no Markdown round-trip). Branding settings ride on
    ``model.export`` for renderers that honor them natively. ``language``
    (Phase 7) localizes the furniture labels. Fail-open: on any error return a
    bare model so a render never breaks."""
    from .model import ExportSettings, markdown_to_model
    md = _apply_branding_md(content, export_settings)
    try:
        model = markdown_to_model(md, title)
    except Exception:  # noqa: BLE001
        model = markdown_to_model("", title)
    if language:
        try:
            model.metadata.language = language
        except Exception:  # noqa: BLE001
            pass
    try:
        from .structure import auto_structure_enabled, enrich
        if auto_structure_enabled():
            model = enrich(model, lang=language)
    except Exception:  # noqa: BLE001
        pass
    model.export = export_settings or ExportSettings()
    return model


def apply_resume_template(content: str, template: str | None) -> str:
    """Phase-4 #21 design layer: re-lay-out resume content through a named,
    ATS-safe template (``classic`` / ``modern`` / ``compact``) before it is
    rendered. Fail-open and identity-preserving — no template, an unknown one,
    or content that isn't a resume returns ``content`` unchanged."""
    if not template:
        return content
    try:
        from .templates import apply_template
        return apply_template(content, template)
    except Exception:  # noqa: BLE001 — design must never break a download
        return content


def render_document(content: str, fmt: str, title: str = "",
                    export_settings=None,
                    template: str | None = None,
                    language: str = "") -> tuple[bytes, str, str]:
    """Return (bytes, media_type, canonical_ext) for the document.

    ``export_settings`` (Phase 7 branding) + Phase-4 structure enrichment are
    applied to EVERY prose format (PDF/DOCX/PPTX/TXT/MD/HTML) via the shared
    DocumentModel, so they render consistently. Data formats (CSV/XLSX/JSON) and
    archives are unaffected.

    ``template`` (Phase-4 #21) re-lays-out RESUME content through a named
    ATS-safe design template first. Default ``None`` → the content is passed
    through untouched, i.e. exactly the legacy behavior.

    ``language`` (Phase 7 localization) localizes the auto-generated furniture
    (Table of Contents / Glossary / Appendix / Table / Figure captions). Blank
    or English → the default English labels (unchanged output)."""
    f = normalize_format(fmt)
    raw = apply_resume_template(content or "", template)
    content = _strip_outer_fence(raw)
    # Phase 1: PDF/DOCX/PPTX render straight from the DocumentModel IR (built +
    # enriched once by `_prose_model`). The legacy Markdown-tuple renderers stay
    # as a fallback behind `model_driven_render`. TXT/MD keep the Markdown path
    # (they ARE Markdown/plain text); HTML consumes the model in its own branch.
    _binary_prose = ("docx", "pptx", "pdf")
    _md_render = _model_driven_enabled()
    # Formats still enriched at the MARKDOWN level: TXT/MD always, and the binary
    # trio only when model-driven rendering is disabled.
    _legacy_prose = ("txt", "md") + (() if _md_render else _binary_prose)
    if f in _legacy_prose:
        content = _enrich_prose_markdown(_apply_branding_md(content,
                                                            export_settings),
                                         title, language)
    if f in _binary_prose and _md_render:
        model = _prose_model(content, title, export_settings, language)
        if f == "docx":
            data = _model_to_docx(model, title)
        elif f == "pptx":
            data = _model_to_pptx(model, title)
        else:
            data = _model_to_pdf(model, title)
        return data, _MEDIA[f], f
    if f == "txt":
        data = _markdown_to_plain(content).encode("utf-8")
    elif f == "md":
        data = content.encode("utf-8")
    elif f == "html":
        # Model-driven (Phase 1): parse into the structured DocumentModel and
        # render HTML from it — the first exporter to consume the IR directly.
        # Phase 4: optionally enrich with a TOC + glossary (config-gated).
        from .model import markdown_to_model, model_to_html
        _model = markdown_to_model(content, title)
        if language:
            try:
                _model.metadata.language = language
            except Exception:  # noqa: BLE001
                pass
        try:
            from .structure import auto_structure_enabled, enrich
            if auto_structure_enabled():
                _model = enrich(_model, lang=language)
        except Exception:  # noqa: BLE001
            pass
        if export_settings is not None:
            _model.export = export_settings
        data = model_to_html(_model).encode("utf-8")
    elif f == "csv":
        data = _to_csv(content).encode("utf-8")
    elif f == "json":
        data = _to_json(content).encode("utf-8")
    elif f == "xlsx":
        data = _to_xlsx(content, title)
    elif f == "docx":
        data = _to_docx(content, title)
    elif f == "pptx":
        data = _to_pptx(content, title)
    elif f == "pdf":
        data = _to_pdf(content, title)
    elif f == "zip":
        # ZIP keeps the ORIGINAL content (fences intact) so every code block
        # becomes a real file in the project archive.
        data = _to_zip(raw, title)
    elif f == "7z":
        # Same project files/structure as the ZIP, LZMA2-compressed.
        data = _to_7z(raw, title)
    elif f == "code" or f in _code_exts():
        # Single source-code file: extract the code from the answer's fenced
        # block and return it raw, with the language's extension.
        from .detect import infer_code_ext

        ext = infer_code_ext(raw) if f == "code" else f
        return _to_code(raw, ext).encode("utf-8"), "text/plain; charset=utf-8", ext
    else:  # pragma: no cover — normalize_format guards this
        raise UnsupportedFormat(fmt)
    return data, _MEDIA[f], f


def _code_exts() -> set[str]:
    from .detect import _CODE_EXTS

    return _CODE_EXTS


def _to_code(content: str, ext: str) -> str:
    """Extract a single source file's text from the answer. Prefers a fenced
    block whose language matches `ext`, else the first fenced block, else the
    whole content."""
    blocks = re.findall(
        r"```([A-Za-z0-9+#.\-]*)[^\n]*\n(.*?)```", content or "", re.DOTALL
    )
    if blocks:
        from .detect import _CODE_LANG

        for lang, body in blocks:
            mapped = _CODE_LANG.get(lang.strip().lower())
            if mapped and mapped == ext:
                return body.strip("\n") + "\n"
        return blocks[0][1].strip("\n") + "\n"
    return (content or "").strip() + "\n"


# --------------------------------------------------------------------------
# Markdown parsing helpers
# --------------------------------------------------------------------------
_SEP_CELL = re.compile(r"^:?-{1,}:?$")


def _strip_outer_fence(text: str) -> str:
    """If the whole content is a single ```fenced``` block, unwrap it."""
    s = text.strip()
    if s.startswith("```") and s.endswith("```"):
        lines = s.splitlines()
        if len(lines) >= 2:
            return "\n".join(lines[1:-1]).strip("\n")
    return text


def _is_table_row(line: str) -> bool:
    s = line.strip()
    return s.startswith("|") and s.count("|") >= 2


def _split_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


def _is_separator(cells: list[str]) -> bool:
    cells = [c for c in cells if c != ""]
    return bool(cells) and all(_SEP_CELL.match(c.replace(" ", "")) for c in cells)


def parse_md_tables(content: str) -> list[list[list[str]]]:
    """Pull every GFM table out of the content as a list of row-lists.

    The `---` separator row is dropped. Non-table lines end the current
    table. Returns [] when there are no tables.
    """
    tables: list[list[list[str]]] = []
    cur: list[list[str]] = []
    for line in content.splitlines():
        if _is_table_row(line):
            cells = _split_row(line)
            if _is_separator(cells):
                continue
            cur.append([_strip_inline(c) for c in cells])
        elif cur:
            tables.append(cur)
            cur = []
    if cur:
        tables.append(cur)
    return tables


def _strip_inline(text: str) -> str:
    """Drop the common inline Markdown markers for plain-text targets."""
    t = text
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"__(.+?)__", r"\1", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)\*", r"\1", t)
    t = re.sub(r"`([^`]+)`", r"\1", t)
    t = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", t)  # links -> text
    return t.strip()


def _markdown_to_plain(content: str) -> str:
    """Readable plain-text rendering: strip fences/markers, keep structure."""
    out: list[str] = []
    for line in content.splitlines():
        s = line.rstrip()
        if s.strip().startswith("```"):
            continue  # drop fence lines, keep the code text
        s = re.sub(r"^#{1,6}\s*", "", s)  # heading markers
        s = re.sub(r"^\s*>\s?", "", s)  # blockquote markers
        out.append(_strip_inline(s) if not _is_table_row(s) else s)
    return "\n".join(out).strip() + "\n"


# Semantic blocks for the docx / pdf renderers.
def parse_blocks(content: str) -> list[tuple]:
    """Parse Markdown into ('kind', payload) blocks.

    Kinds: ('h', level, text), ('p', text), ('bullet', text),
    ('number', text), ('code', text), ('quote', text), ('table', rows).
    """
    blocks: list[tuple] = []
    lines = content.splitlines()
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        s = line.strip()

        # Fenced code block
        if s.startswith("```"):
            lang = s[3:].strip().split(" ")[0].lower() if len(s) > 3 else ""
            i += 1
            code: list[str] = []
            while i < n and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1  # closing fence
            blocks.append(("code", "\n".join(code), lang))
            continue

        # Table
        if _is_table_row(line):
            rows: list[list[str]] = []
            while i < n and _is_table_row(lines[i]):
                cells = _split_row(lines[i])
                if not _is_separator(cells):
                    rows.append([_strip_inline(c) for c in cells])
                i += 1
            if rows:
                blocks.append(("table", rows))
            continue

        if not s:
            i += 1
            continue

        # Standalone image: ![alt](url)  →  ('image', url, alt)
        m = re.match(r"^!\[([^\]]*)\]\(([^)]+)\)\s*$", s)
        if m:
            blocks.append(("image", m.group(2).strip(), m.group(1).strip()))
            i += 1
            continue

        m = re.match(r"^(#{1,6})\s+(.*)$", s)
        if m:
            blocks.append(("h", len(m.group(1)), _strip_inline(m.group(2))))
            i += 1
            continue

        if s.startswith(("- ", "* ", "+ ")):
            blocks.append(("bullet", _strip_inline(s[2:])))
            i += 1
            continue

        m = re.match(r"^\d+[.)]\s+(.*)$", s)
        if m:
            blocks.append(("number", _strip_inline(m.group(1))))
            i += 1
            continue

        if s.startswith(">"):
            blocks.append(("quote", _strip_inline(s.lstrip(">").strip())))
            i += 1
            continue

        blocks.append(("p", _strip_inline(s)))
        i += 1
    return blocks


# --------------------------------------------------------------------------
# Document intelligence (Phase 10, #9) — section numbering + TOC, shared by
# the PDF and DOCX renderers.
# --------------------------------------------------------------------------
_MIN_HEADINGS_FOR_TOC = 3

# Navigational sections that Phase-4 `enrich` injects (or a user writes by hand).
# They must NOT be hierarchically numbered ("2  Table of Contents" is wrong) and
# never appear as their own TOC entry — they ARE the structure, not content.
_STRUCTURAL_HEADINGS = {
    "table of contents", "contents", "glossary", "appendix",
    "list of figures & tables", "list of figures and tables",
    "list of figures", "list of tables", "references", "index",
}


def _is_structural_heading(text: str) -> bool:
    return (text or "").strip().lower() in _STRUCTURAL_HEADINGS


def _structure_active() -> bool:
    """True when Phase-4 enrichment owns document structure (a TOC section, etc.),
    so the legacy renderer-level 'Contents' block must stand down to avoid a
    double table of contents. Fail-open to False (legacy TOC renders)."""
    try:
        from .structure import auto_structure_enabled
        return bool(auto_structure_enabled())
    except Exception:  # noqa: BLE001
        return False


def _doc_flags():
    """(section_numbering, table_of_contents, excel_charts) from config, with
    safe defaults so the generators stay importable without a config."""
    try:
        from app.core.config_loader import cfg
        d = cfg.documents
        return (bool(d.section_numbering), bool(d.table_of_contents),
                bool(d.excel_charts))
    except Exception:  # noqa: BLE001
        return (True, True, True)


def number_headings(blocks: list[tuple]) -> tuple[list[tuple], list[tuple]]:
    """Prepend hierarchical numbers (1, 1.1, 1.1.1) to heading blocks and build
    a table of contents. Returns (numbered_blocks, toc) where toc is a list of
    (level, number, text). No-ops (returns blocks, []) for documents with fewer
    than `_MIN_HEADINGS_FOR_TOC` headings."""
    headings = [b for b in blocks
                if b[0] == "h" and not _is_structural_heading(b[2])]
    if len(headings) < _MIN_HEADINGS_FOR_TOC:
        return blocks, []
    counters = [0] * 6
    out: list[tuple] = []
    toc: list[tuple] = []
    for b in blocks:
        if b[0] == "h" and not _is_structural_heading(b[2]):
            lvl = min(max(int(b[1]), 1), 6)
            counters[lvl - 1] += 1
            for k in range(lvl, 6):
                counters[k] = 0
            number = ".".join(str(counters[i]) for i in range(lvl))
            text = b[2]
            out.append(("h", b[1], f"{number}  {text}"))
            toc.append((lvl, number, text))
        else:
            out.append(b)
    return out, toc


def number_sections(model):
    """Model-level counterpart of :func:`number_headings` (Phase 1): prepend
    hierarchical numbers (1, 1.1, …) to the section headings IN PLACE and return
    ``(model, toc)`` where ``toc`` is a list of ``(level, number, text)``. No-op
    (returns ``model, []``) below ``_MIN_HEADINGS_FOR_TOC`` headings, matching the
    Markdown path exactly."""
    heads = [s for s in model.sections
             if s.heading and not _is_structural_heading(s.heading)]
    if len(heads) < _MIN_HEADINGS_FOR_TOC:
        return model, []
    counters = [0] * 6
    toc: list[tuple] = []
    for s in model.sections:
        if not s.heading or _is_structural_heading(s.heading):
            continue
        lvl = min(max(int(s.level or 1), 1), 6)
        counters[lvl - 1] += 1
        for k in range(lvl, 6):
            counters[k] = 0
        number = ".".join(str(counters[i]) for i in range(lvl))
        toc.append((lvl, number, s.heading))
        s.heading = f"{number}  {s.heading}"
    return model, toc


def _model_stream(model):
    """Linear block stream from a DocumentModel: each section's heading (as a
    :class:`~app.documents.model.Heading`) followed by its blocks — so the DOCX
    and PDF renderers can walk it exactly like the old tuple stream, but off the
    typed IR."""
    from .model import Heading
    for sec in model.sections:
        if sec.heading:
            yield Heading(text=sec.heading, level=sec.level or 1)
        for b in sec.blocks:
            yield b


# --------------------------------------------------------------------------
# CSV / XLSX
# --------------------------------------------------------------------------
def _rows_from_content(content: str) -> list[list[str]]:
    """Best-effort rows for the FIRST table, or CSV-ish / line content."""
    tables = parse_md_tables(content)
    if tables:
        return tables[0]
    # Already comma-separated on most lines? Parse as CSV.
    lines = [ln for ln in content.splitlines() if ln.strip()]
    if lines and sum("," in ln for ln in lines) >= max(1, len(lines) // 2):
        return [next(_csv.reader([ln])) for ln in lines]
    # Fallback: one cell per line.
    return [[_strip_inline(ln)] for ln in lines]


def _to_csv(content: str) -> str:
    buf = io.StringIO()
    writer = _csv.writer(buf)
    tables = parse_md_tables(content)
    rows = tables[0] if tables else _rows_from_content(content)
    for r in rows:
        writer.writerow(r)
    return buf.getvalue()


def _coerce_num(value: str):
    """Convert a numeric-looking cell to int/float (stripping thousands commas
    and a trailing %); otherwise return the original string."""
    s = (value or "").strip()
    if not s:
        return s
    raw = s.rstrip("%").replace(",", "").replace("$", "").strip()
    try:
        if raw and (raw.lstrip("-").isdigit()):
            return int(raw)
        f = float(raw)
        return f
    except (ValueError, TypeError):
        return value


def _numeric_columns(rows: list[list[str]]) -> list[int]:
    """Indices of columns (excluding the header row) that are mostly numeric."""
    if len(rows) < 2:
        return []
    cols = max(len(r) for r in rows)
    numeric: list[int] = []
    for c in range(cols):
        vals = [r[c] for r in rows[1:] if c < len(r) and str(r[c]).strip()]
        if not vals:
            continue
        hits = sum(1 for v in vals if not isinstance(_coerce_num(v), str))
        if hits >= max(1, len(vals) // 2):
            numeric.append(c)
    return numeric


def _write_sheet(ws, rows: list[list[str]], *, charts: bool) -> None:
    """Write a table to a sheet with a bold header, numeric coercion, and
    (when `charts`) a SUM totals row + a bar chart for the first numeric col."""
    from openpyxl.styles import Font

    numeric = _numeric_columns(rows) if charts else []
    for ri, row in enumerate(rows):
        if ri == 0:
            ws.append(row)
            for cell in ws[ws.max_row]:
                cell.font = Font(bold=True)
        else:
            ws.append([
                _coerce_num(v) if (ci in numeric) else v
                for ci, v in enumerate(row)
            ])
    _autosize(ws)
    if not (charts and numeric and ws.max_row >= 3):
        return

    from openpyxl.chart import BarChart, Reference
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    data_rows = ws.max_row
    total_row = data_rows + 1
    ws.cell(row=total_row, column=1, value="Total").font = Font(bold=True)
    for c in numeric:
        col = c + 1
        letter = get_column_letter(col)
        ws.cell(row=total_row, column=col,
                value=f"=SUM({letter}2:{letter}{data_rows})").font = \
            Font(bold=True)

    # Bar chart: first column as categories, first numeric column as values.
    try:
        cat_col = 1
        val_col = numeric[0] + 1
        chart = BarChart()
        chart.title = ws.title
        chart.type = "col"
        chart.style = 10
        data = Reference(ws, min_col=val_col, min_row=1, max_row=data_rows)
        cats = Reference(ws, min_col=cat_col, min_row=2, max_row=data_rows)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.height = 8
        chart.width = 16
        anchor = f"{get_column_letter(ws.max_column + 2)}2"
        ws.add_chart(chart, anchor)
    except Exception:  # noqa: BLE001 — chart is a nicety; never fail the export
        pass


def _to_xlsx(content: str, title: str) -> bytes:
    from openpyxl import Workbook

    _, _, excel_charts = _doc_flags()
    wb = Workbook()
    tables = parse_md_tables(content)
    if not tables:
        ws = wb.active
        ws.title = (title or "Document")[:31] or "Document"
        for r in _rows_from_content(content):
            ws.append(r)
        _autosize(ws)
    else:
        first = True
        for idx, rows in enumerate(tables, start=1):
            ws = wb.active if first else wb.create_sheet()
            ws.title = (f"{title} {idx}" if title else f"Table {idx}")[:31]
            first = False
            _write_sheet(ws, rows, charts=excel_charts)
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _autosize(ws) -> None:
    from openpyxl.utils import get_column_letter

    widths: dict[int, int] = {}
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is None:
                continue
            widths[cell.column] = max(
                widths.get(cell.column, 0), len(str(cell.value))
            )
    for col, w in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = min(60, w + 2)


# --------------------------------------------------------------------------
# JSON — pretty-printed, valid JSON extracted from the answer
# --------------------------------------------------------------------------
def _extract_json(content: str):
    """Best-effort: return a parsed Python object from `content`, or None.

    Tries, in order: the whole content, the first ```json fenced block, any
    fenced block, JSON Lines (one object per line), then the first balanced
    {...} / [...] span found anywhere in the text.
    """
    import json

    s = (content or "").strip()
    if not s:
        return None

    def _try(text: str):
        try:
            return json.loads(text)
        except Exception:  # noqa: BLE001
            return None

    # 1) The whole thing.
    obj = _try(s)
    if obj is not None:
        return obj

    # 2) A fenced block — prefer ```json, else the first fence of any language.
    fences = re.findall(r"```[ \t]*(\w+)?[ \t]*\n(.*?)```", s, re.DOTALL)
    json_first = sorted(
        fences, key=lambda f: 0 if (f[0] or "").lower() == "json" else 1
    )
    for _lang, body in json_first:
        obj = _try(body.strip())
        if obj is not None:
            return obj

    # 3) JSON Lines — each non-empty line is its own object.
    lines = [ln for ln in s.splitlines() if ln.strip()]
    if len(lines) > 1:
        parsed = [_try(ln) for ln in lines]
        if all(p is not None for p in parsed):
            return parsed

    # 4) First balanced {...} or [...] span in the text.
    for opener, closer in (("{", "}"), ("[", "]")):
        start = s.find(opener)
        if start == -1:
            continue
        depth = 0
        in_str = False
        esc = False
        for i in range(start, len(s)):
            ch = s[i]
            if in_str:
                if esc:
                    esc = False
                elif ch == "\\":
                    esc = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
            elif ch == opener:
                depth += 1
            elif ch == closer:
                depth -= 1
                if depth == 0:
                    obj = _try(s[start:i + 1])
                    if obj is not None:
                        return obj
                    break
    return None


def _to_json(content: str) -> str:
    """Return pretty-printed, valid JSON extracted from the answer.

    If no JSON structure can be found, the content is preserved losslessly as
    a single JSON string so the download is always valid JSON.
    """
    import json

    obj = _extract_json(content)
    if obj is None:
        # No structure detected — keep the text, wrapped as valid JSON.
        obj = {"content": (content or "").strip()}
    return json.dumps(obj, indent=2, ensure_ascii=False, sort_keys=False) + "\n"


# --------------------------------------------------------------------------
# ZIP — a full project archive built from the answer's code artifacts
# --------------------------------------------------------------------------
def _safe_zip_name(name: str, used: set[str]) -> str:
    """Sanitize an artifact filename into a safe, unique relative zip path."""
    name = (name or "").strip().replace("\\", "/").lstrip("/")
    # Drop any parent-dir traversal and empty segments.
    parts = [p for p in name.split("/") if p not in ("", ".", "..")]
    safe = "/".join(parts) or "file.txt"
    base = safe
    i = 2
    while safe in used:
        if "." in base.rsplit("/", 1)[-1]:
            stem, ext = base.rsplit(".", 1)
            safe = f"{stem}-{i}.{ext}"
        else:
            safe = f"{base}-{i}"
        i += 1
    used.add(safe)
    return safe


def _parse_dir_tree(content: str) -> list[str]:
    """Extract file paths from an ASCII directory tree in the answer, e.g.

        solar-ui/
        ├─ src/
        │  ├─ App.tsx
        │  └─ index.tsx
        └─ package.json

    Returns full relative paths (e.g. ['src/App.tsx', 'src/index.tsx',
    'package.json']), dropping the top-level project folder. [] if no tree.
    """
    out: list[str] = []
    stack: list[tuple[int, str]] = []  # (indent column, dir name)
    root_seen = False
    for raw in (content or "").splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        m = re.search(r"[├└]", line)
        if m:
            indent = m.start()
            name = line[m.end():].lstrip("─\u2500 ").strip()
        elif (line.strip().endswith("/")
              and not any(c in line for c in "│├└\u2500")
              and re.match(r"^[\w.\-]+/$", line.strip())):
            # A bare root folder line like "solar-ui/".
            indent = -1
            name = line.strip()
            root_seen = True
        else:
            continue
        if not name:
            continue
        # Trim trailing comments after the filename ("main.py   # entry").
        name = re.split(r"\s{2,}|\s#|\s//|\s<!--|\s\u2190|\s<-", name)[0].strip()
        if not name or name in (".", ".."):
            continue
        is_dir = name.endswith("/")
        nm = name.rstrip("/")
        while stack and stack[-1][0] >= indent:
            stack.pop()
        prefix = "/".join(d for _, d in stack)
        full = f"{prefix}/{nm}" if prefix else nm
        if is_dir:
            stack.append((indent, nm))
        elif "." in nm or "/" in full:  # looks like a file
            out.append(full)
    if not out:
        return []
    # Drop the common top-level project folder (e.g. "solar-ui/") so paths are
    # relative to the project root (matching how files were labelled).
    roots = {p.split("/", 1)[0] for p in out if "/" in p}
    if root_seen and len(roots) == 1:
        r = next(iter(roots)) + "/"
        out = [p[len(r):] if p.startswith(r) else p for p in out]
    # De-dup, keep order.
    seen: set[str] = set()
    uniq = []
    for p in out:
        if p and p not in seen:
            seen.add(p)
            uniq.append(p)
    return uniq


def _stub_for(path: str) -> str:
    """Placeholder body for a tree file the model didn't actually write."""
    ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""
    cmt = {
        "py": "#", "rb": "#", "sh": "#", "yml": "#", "yaml": "#", "toml": "#",
        "js": "//", "jsx": "//", "ts": "//", "tsx": "//", "go": "//",
        "java": "//", "kt": "//", "swift": "//", "c": "//", "cpp": "//",
        "cs": "//", "php": "//", "rs": "//", "css": "/*", "scss": "/*",
    }.get(ext, "#")
    close = " */" if cmt == "/*" else ""
    return f"{cmt} {path} — placeholder; not generated in the response.{close}\n"


def _collect_project_files(content: str) -> tuple[str, list[tuple[str, str]]]:
    """Resolve the project's files from the answer's fenced code blocks.

    If the answer shows a directory tree, the file set MATCHES that tree
    exactly (real content where a code block exists, a placeholder otherwise)
    so the archive mirrors the structure the user was shown. Otherwise it's the
    resolved code artifacts. Returns (readme_text, [(path, body), ...]).
    Shared by the ZIP and 7z builders so both archives are identical in layout.
    """
    from ..response_arch.artifacts import split_artifacts

    artifacts = split_artifacts(content or "")
    _snip = re.compile(r"^snippet-\d+\.", re.IGNORECASE)
    by_path: dict[str, str] = {}
    by_base: dict[str, str] = {}
    for a in artifacts:
        if _snip.match(a.filename):
            continue
        norm = (a.filename or "").strip().lstrip("/")
        if not norm:
            continue
        by_path[norm] = a.content or ""
        by_base[norm.rsplit("/", 1)[-1]] = a.content or ""

    tree = _parse_dir_tree(content or "")
    files: list[tuple[str, str]] = []
    if tree:
        for p in tree:
            base = p.rsplit("/", 1)[-1]
            body = by_path.get(p)
            if body is None:
                body = by_base.get(base)
            files.append((p, body if body is not None else _stub_for(p)))
    else:
        real = [(p, c) for p, c in by_path.items()]
        files = real if real else [
            (a.filename, a.content or "") for a in artifacts
        ]
    return (content or "").strip(), files


def _to_zip(content: str, title: str) -> bytes:
    """Build a downloadable project ZIP from the answer's fenced code blocks.

    The archive is built by a script EXECUTED IN THE SANDBOX (same isolation
    as project verification); the in-process zipfile path below is the
    fail-open fallback so a download never breaks."""
    import zipfile

    readme, files = _collect_project_files(content)
    try:
        from app.verify.archive_build import build_archive_sandboxed
        data = build_archive_sandboxed(readme, files, "zip")
        if data:
            return data
    except Exception:  # noqa: BLE001 — fall back to the in-process builder
        pass
    buf = io.BytesIO()
    used: set[str] = set()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        if readme:
            used.add("README.md")
            zf.writestr("README.md", readme)
        for name, body in files:
            fname = _safe_zip_name(name, used)
            zf.writestr(fname, (body or "").rstrip("\n") + "\n")
    return buf.getvalue()


def _to_7z(content: str, title: str) -> bytes:
    """Build a downloadable project .7z — same files/structure as [_to_zip],
    just LZMA2-compressed via py7zr (already a dependency). Falls back to ZIP
    bytes only if py7zr is somehow unavailable, so a download never fails."""
    try:
        import py7zr
    except Exception:  # noqa: BLE001 — keep downloads working without py7zr
        return _to_zip(content, title)

    readme, files = _collect_project_files(content)
    try:
        from app.verify.archive_build import build_archive_sandboxed
        data = build_archive_sandboxed(readme, files, "7z")
        if data:
            return data
    except Exception:  # noqa: BLE001 — fall back to the in-process builder
        pass
    buf = io.BytesIO()
    used: set[str] = set()
    with py7zr.SevenZipFile(buf, "w") as z:
        if readme:
            used.add("README.md")
            z.writestr(readme, "README.md")
        for name, body in files:
            fname = _safe_zip_name(name, used)
            z.writestr((body or "").rstrip("\n") + "\n", fname)
    return buf.getvalue()


# --------------------------------------------------------------------------
# DOCX
# --------------------------------------------------------------------------
def _to_docx(content: str, title: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    blocks = parse_blocks(content)
    sec_num, toc_on, _ = _doc_flags()
    toc: list[tuple] = []
    if sec_num:
        blocks, toc = number_headings(blocks)
    if toc_on and toc and not _structure_active():
        doc.add_heading("Contents", level=1)
        for lvl, number, text in toc:
            p = doc.add_paragraph(f"{number}  {text}")
            p.paragraph_format.left_indent = Pt(12 * max(0, lvl - 1))
        doc.add_page_break()

    for block in blocks:
        kind = block[0]
        if kind == "h":
            doc.add_heading(block[2], level=min(block[1], 4))
        elif kind == "bullet":
            doc.add_paragraph(block[1], style="List Bullet")
        elif kind == "number":
            doc.add_paragraph(block[1], style="List Number")
        elif kind == "quote":
            p = doc.add_paragraph(block[1])
            p.style = "Intense Quote" if "Intense Quote" in [
                s.name for s in doc.styles
            ] else "Quote"
        elif kind == "code":
            lang = block[2] if len(block) > 2 else ""
            if lang == "mermaid":
                cap = doc.add_paragraph()
                run = cap.add_run("Diagram (Mermaid)")
                run.italic = True
                run.font.size = Pt(9)
            for ln in block[1].splitlines() or [""]:
                p = doc.add_paragraph()
                run = p.add_run(ln)
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif kind == "image":
            data = _load_image(block[1])
            if data:
                from docx.shared import Inches

                try:
                    doc.add_picture(io.BytesIO(data), width=Inches(6.0))
                except Exception:  # noqa: BLE001 — unsupported → alt text
                    doc.add_paragraph(f"[image: {block[2] or block[1]}]")
            else:
                doc.add_paragraph(f"[image: {block[2] or block[1]}]")
        elif kind == "table":
            rows = block[1]
            cols = max(len(r) for r in rows)
            t = doc.add_table(rows=0, cols=cols)
            t.style = "Light Grid Accent 1" if "Light Grid Accent 1" in [
                s.name for s in doc.styles
            ] else "Table Grid"
            for ri, row in enumerate(rows):
                cells = t.add_row().cells
                for ci in range(cols):
                    cells[ci].text = row[ci] if ci < len(row) else ""
                if ri == 0:
                    for c in cells:
                        for p in c.paragraphs:
                            for r in p.runs:
                                r.font.bold = True
        else:  # paragraph — render inline bold/italic
            p = doc.add_paragraph()
            _add_inline_runs(p, block[1])

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def _docx_caption(doc, text: str) -> None:
    from docx.shared import Pt
    if not text:
        return
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.italic = True
    run.font.size = Pt(9)


def _model_to_docx(model, title: str) -> bytes:
    """Model-driven DOCX (Phase 1): render straight from the DocumentModel IR
    instead of re-parsing Markdown tuples. Structurally equivalent to
    :func:`_to_docx`, plus native captions (Table N / Figure N) when the Phase-4
    enrichment set them."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    from .model import (
        CodeBlock, Diagram, Heading, Image, ListBlock, Paragraph, Quote, Table,
    )

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    sec_num, toc_on, _ = _doc_flags()
    toc: list[tuple] = []
    if sec_num:
        model, toc = number_sections(model)
    # When Phase-4 enrichment is active it inserts a "Table of Contents" section,
    # so the legacy renderer-level Contents block stands down (no double TOC).
    if toc_on and toc and not _structure_active():
        doc.add_heading("Contents", level=1)
        for lvl, number, text in toc:
            p = doc.add_paragraph(f"{number}  {text}")
            p.paragraph_format.left_indent = Pt(12 * max(0, lvl - 1))
        doc.add_page_break()

    style_names = [s.name for s in doc.styles]
    for b in _model_stream(model):
        if isinstance(b, Heading):
            doc.add_heading(b.text, level=min(b.level, 4))
        elif isinstance(b, ListBlock):
            style = "List Number" if b.ordered else "List Bullet"
            for it in b.items:
                doc.add_paragraph(it, style=style)
        elif isinstance(b, Quote):
            p = doc.add_paragraph(b.text)
            p.style = "Intense Quote" if "Intense Quote" in style_names else "Quote"
        elif isinstance(b, (CodeBlock, Diagram)):
            if isinstance(b, Diagram):
                cap = doc.add_paragraph()
                run = cap.add_run(f"Diagram ({b.diagram_kind.title()})")
                run.italic = True
                run.font.size = Pt(9)
            src = b.source if isinstance(b, Diagram) else b.code
            for ln in src.splitlines() or [""]:
                p = doc.add_paragraph()
                run = p.add_run(ln)
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif isinstance(b, Image):
            data = _load_image(b.url)
            placed = False
            if data:
                from docx.shared import Inches
                try:
                    doc.add_picture(io.BytesIO(data), width=Inches(6.0))
                    placed = True
                except Exception:  # noqa: BLE001 — unsupported → alt text
                    placed = False
            if not placed:
                doc.add_paragraph(f"[image: {b.alt or b.url}]")
            _docx_caption(doc, b.caption)
        elif isinstance(b, Table):
            rows = b.rows
            if rows:
                cols = max(len(r) for r in rows)
                t = doc.add_table(rows=0, cols=cols)
                t.style = ("Light Grid Accent 1"
                           if "Light Grid Accent 1" in style_names
                           else "Table Grid")
                for ri, row in enumerate(rows):
                    cells = t.add_row().cells
                    for ci in range(cols):
                        cells[ci].text = row[ci] if ci < len(row) else ""
                    if ri == 0:
                        for c in cells:
                            for p in c.paragraphs:
                                for r in p.runs:
                                    r.font.bold = True
                _docx_caption(doc, b.caption)
        elif isinstance(b, Paragraph):
            p = doc.add_paragraph()
            _add_inline_runs(p, b.text)

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


_INLINE_RE = re.compile(r"(\*\*.+?\*\*|__.+?__|`[^`]+`|\*.+?\*)")


# --------------------------------------------------------------------------
# PPTX (python-pptx) — Markdown → a slide deck (P2-8, doc-skill parity)
# --------------------------------------------------------------------------
_PPTX_MAX_BULLETS = 11        # bullets per content slide before spilling over
_PPTX_MAX_CODE_LINES = 22     # code lines per slide


def _group_slides(blocks: list[tuple], *, default_title: str) -> list[dict]:
    """Group Markdown blocks into slide dicts.

    A level-1/2 heading starts a new content slide (its text = the slide title);
    bullets/numbers/paragraphs/quotes/deeper-headings become its body items
    (with indent levels). Tables and code blocks become their own slides so the
    layout stays clean. Returns a list of {kind, title, ...}.
    """
    slides: list[dict] = []
    cur: dict | None = None

    def _flush():
        nonlocal cur
        if cur and cur["items"]:
            slides.append(cur)
        cur = None

    def _ensure(title: str):
        nonlocal cur
        if cur is None:
            cur = {"kind": "content", "title": title, "items": []}

    for b in blocks:
        kind = b[0]
        if kind == "h" and int(b[1]) <= 2:
            _flush()
            cur = {"kind": "content", "title": b[2], "items": []}
        elif kind == "h":  # deeper heading → a sub-line
            _ensure(default_title)
            cur["items"].append((0, b[2]))
        elif kind in ("bullet", "number", "quote"):
            _ensure(default_title)
            cur["items"].append((1, b[1]))
        elif kind == "p":
            _ensure(default_title)
            cur["items"].append((0, b[1]))
        elif kind == "table":
            title = (cur or {}).get("title") or "Table"
            _flush()
            slides.append({"kind": "table", "title": title, "rows": b[1]})
        elif kind == "code":
            title = (cur or {}).get("title") or "Code"
            _flush()
            slides.append({"kind": "code", "title": title, "code": b[1]})
        elif kind == "image":
            _ensure(default_title)
            cur["items"].append((0, f"[image: {b[2] or b[1]}]"))
    _flush()
    return slides


def _group_slides_model(model, *, default_title: str) -> list[dict]:
    """Model-driven counterpart of :func:`_group_slides` (Phase 1): group the IR
    blocks into slide dicts off the typed model rather than Markdown tuples."""
    from .model import (
        CodeBlock, Diagram, Heading, Image, ListBlock, Paragraph, Quote, Table,
    )

    slides: list[dict] = []
    cur: dict | None = None

    def _flush():
        nonlocal cur
        if cur and cur["items"]:
            slides.append(cur)
        cur = None

    def _ensure(title: str):
        nonlocal cur
        if cur is None:
            cur = {"kind": "content", "title": title, "items": []}

    for b in _model_stream(model):
        if isinstance(b, Heading) and int(b.level) <= 2:
            _flush()
            cur = {"kind": "content", "title": b.text, "items": []}
        elif isinstance(b, Heading):          # deeper heading → a sub-line
            _ensure(default_title)
            cur["items"].append((0, b.text))
        elif isinstance(b, ListBlock):
            _ensure(default_title)
            for it in b.items:
                cur["items"].append((1, it))
        elif isinstance(b, Quote):
            _ensure(default_title)
            cur["items"].append((1, b.text))
        elif isinstance(b, Paragraph):
            _ensure(default_title)
            cur["items"].append((0, b.text))
        elif isinstance(b, Table):
            title = (cur or {}).get("title") or "Table"
            _flush()
            slides.append({"kind": "table", "title": title, "rows": b.rows})
        elif isinstance(b, (CodeBlock, Diagram)):
            title = (cur or {}).get("title") or "Code"
            _flush()
            src = b.source if isinstance(b, Diagram) else b.code
            slides.append({"kind": "code", "title": title, "code": src})
        elif isinstance(b, Image):
            _ensure(default_title)
            cur["items"].append((0, f"[image: {b.alt or b.url}]"))
    _flush()
    return slides


def _model_to_pptx(model, title: str) -> bytes:
    """Model-driven PPTX (Phase 1): build the deck from the DocumentModel IR."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    sec_num, _toc, _ = _doc_flags()
    if sec_num:
        model, _ = number_sections(model)

    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = "Generated with ZapTheTrick"
            except Exception:  # noqa: BLE001
                pass

    slides = _pptx_split_overflow(
        _group_slides_model(model, default_title=title or "Overview"))
    for s in slides:
        if s["kind"] == "content":
            _pptx_content_slide(prs, s, Pt)
        elif s["kind"] == "table":
            _pptx_table_slide(prs, s, Inches, Pt, sw)
        elif s["kind"] == "code":
            _pptx_code_slide(prs, s, Inches, Pt, sw, sh)

    if not prs.slides:  # never produce an empty deck
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title or "Document"

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pptx_split_overflow(slides: list[dict]) -> list[dict]:
    """Split content slides whose bullet list is too long into '(cont.)' slides."""
    out: list[dict] = []
    for s in slides:
        if s["kind"] != "content" or len(s["items"]) <= _PPTX_MAX_BULLETS:
            out.append(s)
            continue
        items = s["items"]
        first = True
        for i in range(0, len(items), _PPTX_MAX_BULLETS):
            chunk = items[i:i + _PPTX_MAX_BULLETS]
            title = s["title"] if first else f"{s['title']} (cont.)"
            out.append({"kind": "content", "title": title, "items": chunk})
            first = False
    return out


def _to_pptx(content: str, title: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    blocks = parse_blocks(content)
    sec_num, _toc, _ = _doc_flags()
    if sec_num:
        blocks, _ = number_headings(blocks)

    if title:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = "Generated with ZapTheTrick"
            except Exception:  # noqa: BLE001
                pass

    slides = _pptx_split_overflow(
        _group_slides(blocks, default_title=title or "Overview"))
    for s in slides:
        if s["kind"] == "content":
            _pptx_content_slide(prs, s, Pt)
        elif s["kind"] == "table":
            _pptx_table_slide(prs, s, Inches, Pt, sw)
        elif s["kind"] == "code":
            _pptx_code_slide(prs, s, Inches, Pt, sw, sh)

    if not prs.slides:  # never produce an empty deck
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title or "Document"

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pptx_body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def _pptx_content_slide(prs, s: dict, Pt) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = s["title"][:120]
    body = _pptx_body_placeholder(slide)
    if body is None:
        return
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for level, text in s["items"]:
        text = (text or "").strip()
        if not text:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text[:300]
        p.level = min(max(level, 0), 4)
        for r in p.runs:
            r.font.size = Pt(18 if level == 0 else 16)


def _pptx_table_slide(prs, s: dict, Inches, Pt, slide_width) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = s["title"][:120]
    rows = s["rows"] or [[""]]
    nrows = len(rows)
    ncols = max(len(r) for r in rows)
    shape = slide.shapes.add_table(
        nrows, ncols, Inches(0.6), Inches(1.6),
        slide_width - Inches(1.2), Inches(0.4 * nrows))
    table = shape.table
    for ri in range(nrows):
        row = rows[ri]
        for ci in range(ncols):
            cell = table.cell(ri, ci)
            cell.text = (row[ci] if ci < len(row) else "")[:120]
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    if ri == 0:
                        r.font.bold = True


def _pptx_code_slide(prs, s: dict, Inches, Pt, slide_width, slide_height) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = s["title"][:120]
    lines = ((s["code"] or "").splitlines() or [""])[:_PPTX_MAX_CODE_LINES]
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.6), slide_width - Inches(1.2),
        slide_height - Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ln
        for r in p.runs:
            r.font.name = "Consolas"
            r.font.size = Pt(12)


def _add_inline_runs(paragraph, text: str) -> None:
    from docx.shared import Pt

    for part in _INLINE_RE.split(text):
        if not part:
            continue
        if (part.startswith("**") and part.endswith("**")) or (
            part.startswith("__") and part.endswith("__")
        ):
            paragraph.add_run(part[2:-2]).bold = True
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run(part[1:-1])
            run.font.name = "Consolas"
            run.font.size = Pt(9)
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            paragraph.add_run(part[1:-1]).italic = True
        else:
            paragraph.add_run(part)


# --------------------------------------------------------------------------
# PDF (fpdf2 — core fonts are latin-1, so sanitize unicode)
# --------------------------------------------------------------------------
_UNICODE_MAP = {
    "—": "-", "–": "-", "•": "*", "“": '"', "”": '"', "‘": "'", "’": "'",
    "…": "...", "→": "->", "←": "<-", "≥": ">=", "≤": "<=", "×": "x", " ": " ",
}


def _latin1(s: str) -> str:
    for k, v in _UNICODE_MAP.items():
        s = s.replace(k, v)
    return s.encode("latin-1", "replace").decode("latin-1")


# Brand accent (ZapTheTrick purple) used for the title, headings and table head.
_ACCENT = (124, 92, 255)
_INK = (33, 37, 48)
_MUTED = (110, 116, 130)


def _register_fonts(pdf) -> bool:
    """Register DejaVu Unicode fonts as 'Body'(+B/I) and 'Mono', plus an emoji
    fallback. Returns False if the core text fonts are missing (caller degrades
    to the Latin-1 core fonts)."""
    try:
        pdf.add_font("Body", "", os.path.join(_FONT_DIR, "DejaVuSans.ttf"))
        pdf.add_font("Body", "B", os.path.join(_FONT_DIR, "DejaVuSans-Bold.ttf"))
        # No oblique shipped — reuse the regular face for italics.
        pdf.add_font("Body", "I", os.path.join(_FONT_DIR, "DejaVuSans.ttf"))
        pdf.add_font("Mono", "", os.path.join(_FONT_DIR, "DejaVuSansMono.ttf"))
    except Exception:  # noqa: BLE001 — missing fonts → degrade to core fonts
        return False
    # Emoji is optional: a glyph fallback so 🚀/✅/⚠ render instead of blank.
    # exact_match=False so the (regular) emoji font is used even inside bold /
    # italic text (headings, **bold**), where DejaVu lacks the glyph.
    try:
        pdf.add_font("Emoji", "", os.path.join(_FONT_DIR, "Emoji.ttf"))
        pdf.set_fallback_fonts(["Emoji"], exact_match=False)
    except Exception:  # noqa: BLE001
        pass
    return True


def _load_image(url: str) -> bytes | None:
    """Bytes for a markdown image — supports data: URIs and http(s)."""
    import base64
    import urllib.request

    try:
        if url.startswith("data:") and "," in url:
            return base64.b64decode(url.split(",", 1)[1])
        if url.startswith(("http://", "https://")):
            req = urllib.request.Request(
                url, headers={"User-Agent": "Mozilla/5.0 zapthetrick"}
            )
            with urllib.request.urlopen(req, timeout=8) as r:  # noqa: S310
                return r.read()
    except Exception:  # noqa: BLE001
        return None
    return None


def _pdf_image(pdf, url: str, alt: str, epw: float, body: str, txt) -> None:
    data = _load_image(url)
    if not data:
        pdf.set_font(body, "I", 9)
        pdf.set_text_color(*_MUTED)
        pdf.multi_cell(epw, 6, txt(f"[image: {alt or url}]"))
        pdf.set_text_color(0, 0, 0)
        return
    from io import BytesIO

    try:
        from PIL import Image as _PIL

        with _PIL.open(BytesIO(data)) as im:
            iw, ih = im.size
        aspect = ih / max(1, iw)
    except Exception:  # noqa: BLE001
        aspect = 0.62
    w_mm = min(epw, 165.0)
    h_mm = w_mm * aspect
    page_h = pdf.h - pdf.t_margin - pdf.b_margin
    if h_mm > min(215.0, page_h):
        h_mm = min(215.0, page_h)
        w_mm = h_mm / max(0.01, aspect)
    # Start the image on a page with room (fpdf doesn't split an image).
    if h_mm > (pdf.h - pdf.b_margin - pdf.get_y()):
        pdf.add_page()
    pdf.ln(2)
    x = pdf.l_margin + (epw - w_mm) / 2.0
    try:
        pdf.image(BytesIO(data), x=x, w=w_mm, h=h_mm)
    except Exception:  # noqa: BLE001 — unsupported image → show alt
        pdf.set_font(body, "I", 9)
        pdf.set_text_color(*_MUTED)
        pdf.multi_cell(epw, 6, txt(f"[image: {alt or url}]"))
        pdf.set_text_color(0, 0, 0)
    pdf.ln(3)


def _to_pdf(content: str, title: str) -> bytes:
    from fpdf import FPDF

    class _Doc(FPDF):
        foot_font = "Helvetica"

        def footer(self) -> None:  # auto-called per page
            self.set_y(-12)
            try:
                self.set_font(self.foot_font, "", 8)
            except Exception:  # noqa: BLE001
                return
            self.set_draw_color(225, 225, 232)
            self.set_line_width(0.2)
            self.line(self.l_margin, self.get_y(),
                      self.w - self.r_margin, self.get_y())
            self.set_text_color(*_MUTED)
            self.cell(0, 9, f"Page {self.page_no()}", align="C")
            self.set_text_color(0, 0, 0)

    pdf = _Doc(format="A4")
    # Generous, even margins — cleaner look and a buffer so a wide fallback
    # (emoji) glyph can never cross the page edge.
    pdf.set_margins(left=15, top=15, right=15)
    pdf.set_auto_page_break(auto=True, margin=18)

    unicode_ok = _register_fonts(pdf)
    pdf.foot_font = "Body" if unicode_ok else "Helvetica"
    body = "Body" if unicode_ok else "Helvetica"
    mono = "Mono" if unicode_ok else "Courier"
    txt = (lambda s: s) if unicode_ok else _latin1

    pdf.add_page()
    epw = pdf.w - pdf.l_margin - pdf.r_margin

    blocks = parse_blocks(content)
    sec_num, toc_on, _ = _doc_flags()
    toc: list[tuple] = []
    if sec_num:
        blocks, toc = number_headings(blocks)

    # NOTE: every text block uses multi_cell(w=0) — fpdf2 then sizes the cell
    # from the current x to the RIGHT MARGIN, so it's self-correcting and can
    # never run off the page regardless of where a prior block left the cursor.
    if title:
        pdf.set_font(body, "B", 21)
        pdf.set_text_color(*_ACCENT)
        pdf.multi_cell(0, 10, txt(title))
        y = pdf.get_y() + 1.5
        pdf.set_draw_color(*_ACCENT)
        pdf.set_line_width(0.7)
        pdf.line(pdf.l_margin, y, pdf.l_margin + epw, y)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(7)

    # Table of contents (multi-section docs).
    if toc_on and toc and not _structure_active():
        pdf.set_font(body, "B", 14)
        pdf.set_text_color(*_INK)
        pdf.multi_cell(0, 8, txt("Contents"))
        pdf.ln(1)
        pdf.set_font(body, "", 10.5)
        pdf.set_text_color(*_MUTED)
        for lvl, number, text in toc:
            pdf.set_x(pdf.l_margin + 4 * max(0, lvl - 1))
            pdf.multi_cell(0, 5.6, txt(f"{number}  {text}"), wrapmode="CHAR")
        pdf.set_text_color(0, 0, 0)
        pdf.ln(3)

    for block in blocks:
        pdf.set_x(pdf.l_margin)
        kind = block[0]
        if kind == "h":
            lvl = block[1]
            pdf.ln(3 if lvl <= 2 else 2)
            if lvl == 1:
                pdf.set_font(body, "B", 16)
                pdf.set_text_color(*_ACCENT)
                pdf.multi_cell(0, 8, txt(block[2]))
            elif lvl == 2:
                y0 = pdf.get_y()
                pdf.set_fill_color(*_ACCENT)
                pdf.rect(pdf.l_margin, y0 + 0.8, 1.3, 5.4, "F")  # accent bar
                pdf.set_x(pdf.l_margin + 4)
                pdf.set_font(body, "B", 13)
                pdf.set_text_color(*_INK)
                pdf.multi_cell(0, 7, txt(block[2]))
            else:
                pdf.set_font(body, "B", 11.5)
                pdf.set_text_color(*_MUTED)
                pdf.multi_cell(0, 6.5, txt(block[2]))
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1.5)
        elif kind == "bullet":
            pdf.set_font(body, "", 10.5)
            pdf.set_text_color(*_INK)
            pdf.multi_cell(0, 5.8, txt(f"  •  {block[1]}"), wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
        elif kind == "number":
            pdf.set_font(body, "", 10.5)
            pdf.set_text_color(*_INK)
            pdf.multi_cell(0, 5.8, txt(f"  {block[1]}"), wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
        elif kind == "quote":
            y0 = pdf.get_y()
            pdf.set_font(body, "I", 10.5)
            pdf.set_text_color(*_MUTED)
            pdf.set_x(pdf.l_margin + 4)
            pdf.multi_cell(0, 6, txt(block[1]))
            pdf.set_fill_color(*_ACCENT)
            pdf.rect(pdf.l_margin, y0, 1.3, pdf.get_y() - y0, "F")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif kind == "code":
            if len(block) > 2 and block[2] == "mermaid":
                pdf.set_font(body, "I", 9)
                pdf.set_text_color(*_MUTED)
                pdf.multi_cell(0, 5, txt("Diagram (Mermaid)"))
                pdf.set_text_color(0, 0, 0)
            pdf.set_font(mono, "", 7.5)
            pdf.set_fill_color(245, 246, 249)
            pdf.set_text_color(*_INK)
            for ln in block[1].splitlines() or [""]:
                # Reset X to the left margin BEFORE each line so multi_cell
                # always gets the full page width. Without this, a preceding
                # block (e.g. the Mermaid caption) can leave the cursor at the
                # right margin, giving the CHAR-wrap multi_cell ~0 width — an
                # infinite loop in fpdf2.
                pdf.set_x(pdf.l_margin)
                pdf.multi_cell(0, 4.2, txt(ln) or " ", fill=True, wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif kind == "image":
            _pdf_image(pdf, block[1], block[2], epw, body, txt)
        elif kind == "table":
            _pdf_table(pdf, block[1], body, txt)
        else:
            pdf.set_font(body, "", 11)
            pdf.set_text_color(*_INK)
            pdf.multi_cell(0, 6, txt(block[1]), wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)

    return bytes(pdf.output())


def _pdf_heading(pdf, level: int, text: str, body: str, txt) -> None:
    pdf.ln(3 if level <= 2 else 2)
    if level == 1:
        pdf.set_font(body, "B", 16)
        pdf.set_text_color(*_ACCENT)
        pdf.multi_cell(0, 8, txt(text))
    elif level == 2:
        y0 = pdf.get_y()
        pdf.set_fill_color(*_ACCENT)
        pdf.rect(pdf.l_margin, y0 + 0.8, 1.3, 5.4, "F")  # accent bar
        pdf.set_x(pdf.l_margin + 4)
        pdf.set_font(body, "B", 13)
        pdf.set_text_color(*_INK)
        pdf.multi_cell(0, 7, txt(text))
    else:
        pdf.set_font(body, "B", 11.5)
        pdf.set_text_color(*_MUTED)
        pdf.multi_cell(0, 6.5, txt(text))
    pdf.set_text_color(0, 0, 0)
    pdf.ln(1.5)


def _pdf_caption(pdf, text: str, body: str, txt) -> None:
    if not text:
        return
    pdf.set_x(pdf.l_margin)
    pdf.set_font(body, "I", 9)
    pdf.set_text_color(*_MUTED)
    pdf.multi_cell(0, 5, txt(text))
    pdf.set_text_color(0, 0, 0)
    pdf.ln(1)


def _model_to_pdf(model, title: str) -> bytes:
    """Model-driven PDF (Phase 1): render straight from the DocumentModel IR."""
    from fpdf import FPDF

    from .model import (
        CodeBlock, Diagram, Heading, Image, ListBlock, Paragraph, Quote, Table,
    )

    class _Doc(FPDF):
        foot_font = "Helvetica"

        def footer(self) -> None:  # auto-called per page
            self.set_y(-12)
            try:
                self.set_font(self.foot_font, "", 8)
            except Exception:  # noqa: BLE001
                return
            self.set_draw_color(225, 225, 232)
            self.set_line_width(0.2)
            self.line(self.l_margin, self.get_y(),
                      self.w - self.r_margin, self.get_y())
            self.set_text_color(*_MUTED)
            self.cell(0, 9, f"Page {self.page_no()}", align="C")
            self.set_text_color(0, 0, 0)

    pdf = _Doc(format="A4")
    pdf.set_margins(left=15, top=15, right=15)
    pdf.set_auto_page_break(auto=True, margin=18)

    unicode_ok = _register_fonts(pdf)
    pdf.foot_font = "Body" if unicode_ok else "Helvetica"
    body = "Body" if unicode_ok else "Helvetica"
    mono = "Mono" if unicode_ok else "Courier"
    txt = (lambda s: s) if unicode_ok else _latin1

    pdf.add_page()
    epw = pdf.w - pdf.l_margin - pdf.r_margin

    sec_num, toc_on, _ = _doc_flags()
    toc: list[tuple] = []
    if sec_num:
        model, toc = number_sections(model)

    if title:
        pdf.set_font(body, "B", 21)
        pdf.set_text_color(*_ACCENT)
        pdf.multi_cell(0, 10, txt(title))
        y = pdf.get_y() + 1.5
        pdf.set_draw_color(*_ACCENT)
        pdf.set_line_width(0.7)
        pdf.line(pdf.l_margin, y, pdf.l_margin + epw, y)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(7)

    if toc_on and toc and not _structure_active():
        pdf.set_font(body, "B", 14)
        pdf.set_text_color(*_INK)
        pdf.multi_cell(0, 8, txt("Contents"))
        pdf.ln(1)
        pdf.set_font(body, "", 10.5)
        pdf.set_text_color(*_MUTED)
        for lvl, number, text in toc:
            pdf.set_x(pdf.l_margin + 4 * max(0, lvl - 1))
            pdf.multi_cell(0, 5.6, txt(f"{number}  {text}"), wrapmode="CHAR")
        pdf.set_text_color(0, 0, 0)
        pdf.ln(3)

    for b in _model_stream(model):
        pdf.set_x(pdf.l_margin)
        if isinstance(b, Heading):
            _pdf_heading(pdf, b.level, b.text, body, txt)
        elif isinstance(b, ListBlock):
            pdf.set_font(body, "", 10.5)
            pdf.set_text_color(*_INK)
            for it in b.items:
                pdf.set_x(pdf.l_margin)
                prefix = "  " if b.ordered else "  •  "
                pdf.multi_cell(0, 5.8, txt(f"{prefix}{it}"), wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
        elif isinstance(b, Quote):
            y0 = pdf.get_y()
            pdf.set_font(body, "I", 10.5)
            pdf.set_text_color(*_MUTED)
            pdf.set_x(pdf.l_margin + 4)
            pdf.multi_cell(0, 6, txt(b.text))
            pdf.set_fill_color(*_ACCENT)
            pdf.rect(pdf.l_margin, y0, 1.3, pdf.get_y() - y0, "F")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif isinstance(b, (CodeBlock, Diagram)):
            if isinstance(b, Diagram):
                pdf.set_font(body, "I", 9)
                pdf.set_text_color(*_MUTED)
                pdf.multi_cell(0, 5, txt(f"Diagram ({b.diagram_kind.title()})"))
                pdf.set_text_color(0, 0, 0)
            src = b.source if isinstance(b, Diagram) else b.code
            pdf.set_font(mono, "", 7.5)
            pdf.set_fill_color(245, 246, 249)
            pdf.set_text_color(*_INK)
            for ln in src.splitlines() or [""]:
                pdf.set_x(pdf.l_margin)
                pdf.multi_cell(0, 4.2, txt(ln) or " ", fill=True, wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)
        elif isinstance(b, Image):
            _pdf_image(pdf, b.url, b.alt, epw, body, txt)
            _pdf_caption(pdf, b.caption, body, txt)
        elif isinstance(b, Table):
            _pdf_table(pdf, b.rows, body, txt)
            _pdf_caption(pdf, b.caption, body, txt)
        elif isinstance(b, Paragraph):
            pdf.set_font(body, "", 11)
            pdf.set_text_color(*_INK)
            pdf.multi_cell(0, 6, txt(b.text), wrapmode="CHAR")
            pdf.set_text_color(0, 0, 0)
            pdf.ln(1)

    return bytes(pdf.output())


def _pdf_table(pdf, rows: list[list[str]], font: str, txt) -> None:
    """Native fpdf2 table: wraps long cell text, accent header, zebra rows."""
    if not rows:
        return
    from fpdf.fonts import FontFace

    cols = max(len(r) for r in rows)
    norm = [
        [txt(r[i]) if i < len(r) else "" for i in range(cols)] for r in rows
    ]
    pdf.ln(1)
    pdf.set_font(font, "", 8.5)
    pdf.set_draw_color(222, 224, 232)
    with pdf.table(
        text_align="LEFT",
        first_row_as_headings=True,
        headings_style=FontFace(
            emphasis="BOLD", color=(255, 255, 255), fill_color=_ACCENT
        ),
        line_height=pdf.font_size * 1.7,
        borders_layout="MINIMAL",
        cell_fill_color=(245, 245, 250),
        cell_fill_mode="ROWS",
        wrapmode="CHAR",
    ) as table:
        for r in norm:
            trow = table.row()
            for cell in r:
                trow.cell(cell)
    pdf.set_draw_color(0, 0, 0)
    pdf.ln(2)


__all__ = [
    "render_document",
    "apply_resume_template",
    "normalize_format",
    "media_type",
    "SUPPORTED_FORMATS",
    "UnsupportedFormat",
    "parse_md_tables",
    "parse_blocks",
]
