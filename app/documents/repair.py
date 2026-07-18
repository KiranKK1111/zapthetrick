"""
In-place document repair & beautify (roadmap Phase 4 / feature #6).

Unlike `transform.py` — which extracts an upload to TEXT and re-renders (lossy
for binary layouts) — this repairs a document IN ITS NATIVE FORMAT so structure
and content survive:

  xlsx : openpyxl load → beautify (bold+frozen header, sensible column widths,
         number formats) + optional automation (SUM totals) → save
  docx : python-docx load → collapse blank paragraphs, trim trailing spaces,
         apply a consistent base font → save
  pptx : python-pptx load → enable word-wrap, floor tiny fonts → save
  pdf  : PyMuPDF rewrite with garbage-collection + clean (fixes many broken /
         bloated PDFs) → save
  md   : deterministic markdown beautify (heading spacing, list markers,
         trailing whitespace, blank-line collapse, fenced-block balance)
  code : format (black/prettier/gofmt via polyglot) + lint findings

Every path is FAIL-OPEN: an unavailable library or a corrupt file yields a
report with `ok=False` + a reason, never an exception. Binary formats return new
bytes in `.data`; text formats (md/code/text) return the result in `.text`.
"""
from __future__ import annotations

import io
import re
from dataclasses import dataclass, field

_XLSX = {"xlsx", "xlsm"}
_DOCX = {"docx"}
_PPTX = {"pptx"}
_PDF = {"pdf"}
_MD = {"md", "markdown"}
_CODE = {
    "py": "python", "js": "javascript", "ts": "typescript", "go": "go",
    "rs": "rust", "java": "java", "c": "c", "cpp": "cpp", "cc": "cpp",
    "rb": "ruby", "sh": "bash", "php": "php", "cs": "csharp",
}


@dataclass
class RepairResult:
    ok: bool
    kind: str                       # xlsx|docx|pptx|pdf|md|code|text
    data: bytes = b""               # repaired bytes (binary formats)
    text: str = ""                  # beautified text (md/code/text)
    changed: bool = False
    changes: list[str] = field(default_factory=list)
    lint: list[dict] = field(default_factory=list)
    reason: str = ""

    def to_dict(self) -> dict:
        return {"ok": self.ok, "kind": self.kind, "changed": self.changed,
                "changes": self.changes, "lint": self.lint,
                "reason": self.reason,
                "bytes": len(self.data), "chars": len(self.text)}


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in (filename or "") else ""


# --------------------------------------------------------------------------- #
# Markdown
# --------------------------------------------------------------------------- #
def beautify_markdown(text: str) -> tuple[str, list[str]]:
    """Deterministic markdown cleanup. Never raises."""
    changes: list[str] = []
    try:
        src = (text or "").replace("\r\n", "\n").replace("\r", "\n")
        if src != (text or ""):
            changes.append("normalized line endings")
        lines = src.split("\n")
        out: list[str] = []
        in_fence = False
        fences = 0
        for raw in lines:
            line = raw
            stripped = line.strip()
            if stripped.startswith("```"):
                fences += 1
                in_fence = not in_fence
                out.append(line.rstrip())
                continue
            if in_fence:
                out.append(line.rstrip("\n"))   # keep code indentation as-is
                continue
            # trailing whitespace
            line = line.rstrip()
            # "#Heading" -> "# Heading"
            m = re.match(r"^(#{1,6})([^#\s])", line)
            if m:
                line = f"{m.group(1)} {line[len(m.group(1)):].lstrip()}"
                changes.append("added space after heading marker")
            # unordered list markers "*"/"+" -> "-"
            m = re.match(r"^(\s*)[*+](\s+)", line)
            if m:
                line = f"{m.group(1)}-{m.group(2)}{line[m.end():]}"
            out.append(line)
        # blank line around headings + collapse 3+ blanks to 1
        norm: list[str] = []
        for i, line in enumerate(out):
            is_head = bool(re.match(r"^#{1,6}\s", line))
            if is_head and norm and norm[-1].strip() != "":
                norm.append("")
                changes.append("blank line before heading")
            norm.append(line)
            if is_head and i + 1 < len(out) and out[i + 1].strip() != "":
                norm.append("")
        collapsed: list[str] = []
        blanks = 0
        for line in norm:
            if line.strip() == "":
                blanks += 1
                if blanks <= 1:
                    collapsed.append("")
            else:
                blanks = 0
                collapsed.append(line)
        result = "\n".join(collapsed).strip() + "\n"
        if fences % 2 != 0:
            changes.append("WARNING: unbalanced ``` code fence")
        if result != src and "beautified" not in changes:
            changes.append("beautified markdown")
        return result, changes
    except Exception as exc:  # noqa: BLE001
        return text or "", [f"error: {exc}"]


# --------------------------------------------------------------------------- #
# Excel (xlsx)
# --------------------------------------------------------------------------- #
def _numeric_columns(ws) -> list[int]:
    """1-based indices of columns whose DATA rows are mostly numeric."""
    cols: list[int] = []
    for col in range(1, ws.max_column + 1):
        nums = sum(1 for row in range(2, ws.max_row + 1)
                   if isinstance(ws.cell(row=row, column=col).value, (int, float)))
        if nums >= 2:
            cols.append(col)
    return cols


def repair_xlsx(data: bytes, *, automate: bool = False) -> RepairResult:
    """Beautify an .xlsx in place (openpyxl): bold+frozen header, sensible column
    widths, number formats, and an auto-filter. `automate` additionally adds SUM
    and AVERAGE totals rows, a color-scale heat map on numeric columns, and a bar
    chart of the numeric data. Fail-open (per feature)."""
    r = RepairResult(ok=False, kind="xlsx")
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import Alignment, Font
        from openpyxl.utils import get_column_letter
        wb = load_workbook(io.BytesIO(data))
        for ws in wb.worksheets:
            if ws.max_row < 1 or ws.max_column < 1:
                continue
            # --- Beautify (always) -----------------------------------------
            for cell in ws[1]:
                cell.font = Font(bold=True)
                cell.alignment = Alignment(vertical="center")
            ws.freeze_panes = "A2"
            r.changes.append(f"{ws.title}: bold+frozen header")
            numeric = _numeric_columns(ws)
            for col in range(1, ws.max_column + 1):
                width = 8
                for row in range(1, min(ws.max_row, 200) + 1):
                    v = ws.cell(row=row, column=col).value
                    if v is not None:
                        width = max(width, min(60, len(str(v)) + 2))
                ws.column_dimensions[get_column_letter(col)].width = width
                # Thousands/decimal number format on numeric data cells.
                if col in numeric:
                    for row in range(2, ws.max_row + 1):
                        c = ws.cell(row=row, column=col)
                        if isinstance(c.value, float):
                            c.number_format = "#,##0.00"
                        elif isinstance(c.value, int):
                            c.number_format = "#,##0"
            if ws.max_row >= 2 and ws.max_column >= 1:
                try:
                    ws.auto_filter.ref = ws.dimensions
                    r.changes.append(f"{ws.title}: enabled auto-filter")
                except Exception:  # noqa: BLE001
                    pass
            # --- Automation (opt-in) ---------------------------------------
            if automate and ws.max_row >= 2 and numeric:
                last = ws.max_row
                sum_row, avg_row = last + 1, last + 2
                for col in numeric:
                    cl = get_column_letter(col)
                    sc = ws.cell(row=sum_row, column=col)
                    sc.value = f"=SUM({cl}2:{cl}{last})"
                    sc.font = Font(bold=True)
                    ac = ws.cell(row=avg_row, column=col)
                    ac.value = f"=AVERAGE({cl}2:{cl}{last})"
                    ac.number_format = "#,##0.00"
                ws.cell(row=sum_row, column=1).value = "TOTAL"
                ws.cell(row=sum_row, column=1).font = Font(bold=True)
                ws.cell(row=avg_row, column=1).value = "AVERAGE"
                ws.cell(row=avg_row, column=1).font = Font(bold=True)
                r.changes.append(f"{ws.title}: added SUM + AVERAGE rows")
                # Color-scale heat map over the numeric data block.
                try:
                    from openpyxl.formatting.rule import ColorScaleRule
                    rule = ColorScaleRule(
                        start_type="min", start_color="FFF8696B",
                        mid_type="percentile", mid_value=50, mid_color="FFFFEB84",
                        end_type="max", end_color="FF63BE7B")
                    for col in numeric:
                        cl = get_column_letter(col)
                        ws.conditional_formatting.add(
                            f"{cl}2:{cl}{last}", rule)
                    r.changes.append(f"{ws.title}: color-scale heat map")
                except Exception:  # noqa: BLE001
                    pass
                # Bar chart: numeric columns vs the first column as categories.
                try:
                    from openpyxl.chart import BarChart, Reference
                    chart = BarChart()
                    chart.title = f"{ws.title} — overview"
                    chart.type = "col"
                    first_num = numeric[0]
                    dref = Reference(ws, min_col=numeric[0],
                                     max_col=numeric[-1], min_row=1, max_row=last)
                    cats = Reference(ws, min_col=1, min_row=2, max_row=last)
                    chart.add_data(dref, titles_from_data=True)
                    chart.set_categories(cats)
                    anchor = f"{get_column_letter(ws.max_column + 2)}2"
                    ws.add_chart(chart, anchor)
                    r.changes.append(f"{ws.title}: bar chart")
                except Exception:  # noqa: BLE001
                    pass
        buf = io.BytesIO()
        wb.save(buf)
        r.data = buf.getvalue()
        r.ok = True
        r.changed = bool(r.changes)
        return r
    except Exception as exc:  # noqa: BLE001
        r.reason = f"xlsx repair failed: {exc}"
        return r


# --------------------------------------------------------------------------- #
# Word (docx)
# --------------------------------------------------------------------------- #
def repair_docx(data: bytes, *, base_font: str = "Calibri") -> RepairResult:
    """Beautify a .docx in place (python-docx): drop consecutive blank
    paragraphs, trim trailing whitespace, apply a consistent base font.
    Fail-open."""
    r = RepairResult(ok=False, kind="docx")
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        # Consistent base font on the Normal style.
        try:
            doc.styles["Normal"].font.name = base_font
            r.changes.append(f"base font → {base_font}")
        except Exception:  # noqa: BLE001
            pass
        removed = 0
        prev_blank = False
        for para in list(doc.paragraphs):
            text = para.text
            for run in para.runs:
                if run.text and run.text != run.text.rstrip():
                    run.text = run.text.rstrip()
            is_blank = not text.strip()
            if is_blank and prev_blank:
                el = para._element
                el.getparent().remove(el)
                removed += 1
                continue
            prev_blank = is_blank
        if removed:
            r.changes.append(f"removed {removed} blank paragraph(s)")
        buf = io.BytesIO()
        doc.save(buf)
        r.data = buf.getvalue()
        r.ok = True
        r.changed = bool(r.changes)
        return r
    except Exception as exc:  # noqa: BLE001
        r.reason = f"docx repair failed: {exc}"
        return r


# --------------------------------------------------------------------------- #
# PowerPoint (pptx)
# --------------------------------------------------------------------------- #
def repair_pptx(data: bytes, *, min_font_pt: int = 12) -> RepairResult:
    """Beautify a .pptx (python-pptx): enable word-wrap and floor tiny body
    fonts so text doesn't overflow/underread. Conservative + fail-open."""
    r = RepairResult(ok=False, kind="pptx")
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(io.BytesIO(data))
        wrapped = bumped = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                if tf.word_wrap is not True:
                    tf.word_wrap = True
                    wrapped += 1
                for para in tf.paragraphs:
                    for run in para.runs:
                        sz = run.font.size
                        if sz is not None and sz < Pt(min_font_pt):
                            run.font.size = Pt(min_font_pt)
                            bumped += 1
        if wrapped:
            r.changes.append(f"enabled word-wrap on {wrapped} text box(es)")
        if bumped:
            r.changes.append(f"raised {bumped} run(s) to ≥{min_font_pt}pt")
        buf = io.BytesIO()
        prs.save(buf)
        r.data = buf.getvalue()
        r.ok = True
        r.changed = bool(r.changes)
        return r
    except Exception as exc:  # noqa: BLE001
        r.reason = f"pptx repair failed: {exc}"
        return r


# --------------------------------------------------------------------------- #
# PDF (PyMuPDF)
# --------------------------------------------------------------------------- #
def _ocr_page(page) -> int:
    """Overlay an INVISIBLE OCR text layer on a scanned page so it's searchable/
    selectable. Returns the number of words added (0 if OCR unavailable). Needs
    Tesseract; fail-open."""
    try:
        import fitz  # noqa: F401
        tp = page.get_textpage_ocr(flags=0, language="eng", dpi=200, full=True)
        words = page.get_text("words", textpage=tp) or []
        added = 0
        for w in words:
            x0, y0, x1, y1, text = w[0], w[1], w[2], w[3], w[4]
            if not str(text).strip():
                continue
            size = max(4.0, min(40.0, (y1 - y0) * 0.9))
            # render_mode=3 → invisible text (searchable layer over the image).
            page.insert_text((x0, y1), str(text), fontsize=size,
                             render_mode=3)
            added += 1
        return added
    except Exception:  # noqa: BLE001 — no Tesseract / OCR failure → skip
        return 0


def repair_pdf(data: bytes, *, ocr: bool = True) -> RepairResult:
    """Repair/optimize a PDF (PyMuPDF): rewrite with garbage-collection + clean
    (fixes many broken/bloated PDFs, drops orphaned objects). When `ocr` is on,
    scanned/image-only pages get a searchable OCR text layer (needs Tesseract;
    reported as "needs OCR" when it isn't available). Fail-open."""
    r = RepairResult(ok=False, kind="pdf")
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=data, filetype="pdf")
        n = doc.page_count
        # Detect image-only (scanned) pages: images present, no extractable text.
        scanned = []
        for i, page in enumerate(doc):
            try:
                has_text = bool(page.get_text().strip())
                has_img = bool(page.get_images())
                if has_img and not has_text:
                    scanned.append(i)
            except Exception:  # noqa: BLE001
                pass
        ocr_pages = 0
        if ocr and scanned:
            for i in scanned:
                if _ocr_page(doc[i]) > 0:
                    ocr_pages += 1
            if ocr_pages:
                r.changes.append(
                    f"added OCR text layer to {ocr_pages} scanned page(s)")
            else:
                r.changes.append(
                    f"{len(scanned)} scanned page(s) need OCR "
                    "(install Tesseract for a searchable text layer)")
        buf = doc.tobytes(garbage=4, deflate=True, clean=True)
        doc.close()
        r.data = buf
        r.ok = True
        r.changed = len(buf) != len(data) or ocr_pages > 0
        saved = len(data) - len(buf)
        r.changes.append(
            f"rewrote {n} page(s); "
            + (f"saved {saved} bytes" if saved > 0 else "cleaned structure"))
        return r
    except Exception as exc:  # noqa: BLE001
        r.reason = f"pdf repair failed: {exc}"
        return r


# --------------------------------------------------------------------------- #
# Source code (format + lint)
# --------------------------------------------------------------------------- #
async def repair_code(code: str, language: str) -> RepairResult:
    """Repair code: auto-FIX safe lint issues (ruff/eslint --fix), FORMAT it
    (black/prettier/gofmt), then report any lint findings that remain."""
    r = RepairResult(ok=True, kind="code", text=code or "")
    # 1) auto-fix safe lint issues.
    try:
        from app.polyglot.linters import fix_code
        fixed = await fix_code(language, r.text)
        if fixed and fixed.strip() and fixed != r.text:
            r.text = fixed
            r.changed = True
            r.changes.append(f"auto-fixed lint ({language})")
    except Exception:  # noqa: BLE001
        pass
    # 2) format.
    try:
        from app.polyglot.formatters import format_code
        formatted = await format_code(language, r.text)
        if formatted and formatted.strip() and formatted != r.text:
            r.text = formatted
            r.changed = True
            r.changes.append(f"formatted ({language})")
    except Exception:  # noqa: BLE001
        pass
    # 3) report what's left.
    try:
        from app.polyglot.linters import lint_code
        findings = await lint_code(language, r.text)
        r.lint = [f.to_dict() if hasattr(f, "to_dict") else dict(f)
                  for f in findings]
        if r.lint:
            r.changes.append(f"{len(r.lint)} lint finding(s) remain")
    except Exception:  # noqa: BLE001
        pass
    return r


# --------------------------------------------------------------------------- #
# Dispatcher
# --------------------------------------------------------------------------- #
async def repair_document(data: bytes, filename: str, *,
                          automate: bool = False) -> RepairResult:
    """Repair/beautify a document by its type (binary → native in-place repair;
    text/code → beautify). Never raises."""
    ext = _ext(filename)
    try:
        if ext in _XLSX:
            return repair_xlsx(data, automate=automate)
        if ext in _DOCX:
            return repair_docx(data)
        if ext in _PPTX:
            return repair_pptx(data)
        if ext in _PDF:
            return repair_pdf(data)
        # Text formats.
        text = data.decode("utf-8", errors="replace") if data else ""
        if ext in _MD:
            out, changes = beautify_markdown(text)
            return RepairResult(ok=True, kind="md", text=out,
                                changed=out != text, changes=changes)
        if ext in _CODE:
            return await repair_code(text, _CODE[ext])
        # Unknown → treat as markdown-ish text beautify (safe).
        out, changes = beautify_markdown(text)
        return RepairResult(ok=True, kind="text", text=out,
                            changed=out != text, changes=changes)
    except Exception as exc:  # noqa: BLE001
        return RepairResult(ok=False, kind=ext or "unknown",
                            reason=f"repair failed: {exc}")


__all__ = ["RepairResult", "repair_document", "repair_xlsx", "repair_docx",
           "repair_pptx", "repair_pdf", "repair_code", "beautify_markdown"]
