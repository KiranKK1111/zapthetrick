"""Sandbox parse-verification for rendered documents (docx/xlsx/pptx/pdf).

The project archive already runs through the sandbox verify loop
(`project_loop.verify_and_repair_archive`); single documents only had the
in-process structural validators. This closes that gap: the rendered bytes are
written to an ephemeral workspace and RE-OPENED by their own parser library
inside the layered sandbox (bwrap → rlimit → subprocess), exactly the
"generate, then actually execute a script against the artifact" loop Claude
uses. A document that its own reader cannot open never ships silently — the
result is folded into the X-Artifact-Validation meta the FE badge reads.

Fail-open by design: sandbox unavailable / driver error → {"status": "skipped"}
and delivery proceeds (the in-process validators already ran).
"""
from __future__ import annotations

import asyncio
import logging
import os
import sys
import tempfile

log = logging.getLogger(__name__)

# Per-format EXTRACTORS: open artifact.<ext> and pull ALL its rendered text into
# a `txt` variable. This forces a real parse AND lets the shared coverage tail
# check the content actually made it into the file (not just "does it open").
_EXTRACT: dict[str, str] = {
    "docx": (
        "from docx import Document\n"
        "d = Document('artifact.docx')\n"
        "txt = ' '.join(p.text for p in d.paragraphs)\n"
        "for t in d.tables:\n"
        "    for row in t.rows:\n"
        "        for c in row.cells:\n"
        "            txt += ' ' + c.text\n"
    ),
    "xlsx": (
        "from openpyxl import load_workbook\n"
        "wb = load_workbook('artifact.xlsx')\n"
        "txt = ' '.join(str(c.value) for ws in wb.worksheets "
        "for row in ws.iter_rows() for c in row if c.value is not None)\n"
    ),
    "pptx": (
        "from pptx import Presentation\n"
        "p = Presentation('artifact.pptx')\n"
        "txt = ' '.join(sh.text for sl in p.slides for sh in sl.shapes "
        "if sh.has_text_frame)\n"
    ),
    "pdf": (
        "import fitz\n"
        "doc = fitz.open('artifact.pdf')\n"
        "txt = ' '.join(pg.get_text() for pg in doc)\n"
    ),
    # Structured text formats: a real parse (stdlib). json is structural only
    # (no coverage — key order/whitespace differ); csv gets coverage.
    "json": (
        "import json\n"
        "with open('artifact.json', encoding='utf-8') as f:\n"
        "    json.load(f)\n"
        "txt = ''\n"
    ),
    "csv": (
        "import csv\n"
        "with open('artifact.csv', newline='', encoding='utf-8') as f:\n"
        "    rows = list(csv.reader(f))\n"
        "assert rows\n"
        "txt = ' '.join(' '.join(r) for r in rows)\n"
    ),
}

# Shared tail: FUNCTIONAL check beyond "opens". If an expected-content file is
# present, verify the rendered document actually CONTAINS a good fraction of the
# source's significant words — catches a render that silently truncated, dropped
# a section, or came out empty/garbled. Only enforced when there are enough
# expected tokens to be meaningful (short docs skip it). `raise SystemExit(msg)`
# fails the run with the reason on stderr, which the caller reports as "failed".
_COVERAGE_TAIL = (
    "import re, os\n"
    "COV_MIN = {cov_min}\n"
    "have = set(re.findall(r'[a-z0-9]{{3,}}', (txt or '').lower()))\n"
    "want = []\n"
    "if os.path.exists('expect.txt'):\n"
    "    with open('expect.txt', encoding='utf-8') as _f:\n"
    "        want = re.findall(r'[a-z0-9]{{3,}}', _f.read().lower())\n"
    "wset = set(want)\n"
    "if len(wset) >= 25:\n"
    "    cov = sum(1 for w in wset if w in have) / len(wset)\n"
    "    if cov < COV_MIN:\n"
    "        raise SystemExit('CONTENT_COVERAGE %.0f%% < %.0f%% — the render "
    "dropped content' % (cov * 100, COV_MIN * 100))\n"
    "    print('OK cov=%.2f' % cov)\n"
    "else:\n"
    "    print('OK')\n"
)


async def verify_document_bytes(
    data: bytes, ext: str, expect_text: str | None = None) -> dict | None:
    """Open `data` as `ext` inside the sandbox and, when `expect_text` (the
    source content) is given, verify the rendered document actually CONTAINS
    that content (functional check, not just "does it parse"). Returns a small
    meta dict ({"status": "verified"|"failed"|"skipped", ...}) or None for
    formats without a driver (txt/md are plain text — nothing to parse)."""
    extract = _EXTRACT.get((ext or "").lower())
    if extract is None:
        return None
    try:
        cov_min = _coverage_min()
        driver = extract + _COVERAGE_TAIL.format(cov_min=cov_min)
        return await asyncio.to_thread(
            _run, data, ext.lower(), driver, expect_text or "")
    except Exception as exc:  # noqa: BLE001 — never block delivery
        log.debug("doc sandbox verify skipped: %s", exc)
        return {"status": "skipped", "reason": str(exc)[:120]}


def _coverage_min() -> float:
    try:
        from app.core.config_loader import cfg
        return float(getattr(cfg.documents, "verify_coverage_min", 0.6) or 0.6)
    except Exception:  # noqa: BLE001
        return 0.6


def _run(data: bytes, ext: str, driver: str, expect_text: str) -> dict:
    from app.sandbox.executor import SandboxLimits, run_command
    with tempfile.TemporaryDirectory(prefix="docverify-") as tmp:
        with open(os.path.join(tmp, f"artifact.{ext}"), "wb") as f:
            f.write(data)
        with open(os.path.join(tmp, "check.py"), "w", encoding="utf-8") as f:
            f.write(driver)
        if expect_text.strip():
            with open(os.path.join(tmp, "expect.txt"), "w",
                      encoding="utf-8") as f:
                f.write(expect_text)
        limits = SandboxLimits.from_config()
        limits.timeout_s = min(limits.timeout_s, 15)
        res = run_command([sys.executable, "check.py"],
                          workdir=tmp, limits=limits)
    if res.status == "ok":
        _cov = (res.stdout or "").strip()
        return {"status": "verified", "backend": res.backend,
                "ms": res.duration_ms,
                **({"coverage": _cov} if "cov=" in _cov else {})}
    if res.status in ("unavailable", "timeout"):
        return {"status": "skipped", "reason": res.status,
                "backend": res.backend}
    return {"status": "failed", "backend": res.backend,
            "detail": (res.stderr or res.reason or "")[-200:]}
